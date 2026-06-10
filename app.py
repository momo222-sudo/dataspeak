import os, sqlite3, secrets, json, io
from datetime import datetime, date
from flask import Flask, request, jsonify, send_from_directory, redirect, send_file
import anthropic
from dotenv import load_dotenv

load_dotenv()

app = Flask(__name__, static_folder='static')

ANTHROPIC_KEY        = os.environ.get('ANTHROPIC_API_KEY', '')
STRIPE_SECRET_KEY    = os.environ.get('STRIPE_SECRET_KEY', '')
STRIPE_WEBHOOK_SECRET= os.environ.get('STRIPE_WEBHOOK_SECRET', '')
STRIPE_PRICE_ID      = os.environ.get('STRIPE_PRICE_ID', '')
APP_URL              = os.environ.get('APP_URL', 'https://dataspeak-vydp.onrender.com')
DB_PATH              = 'users.db'

FREE_LIMIT = 10   # analyses per month on free plan

# ── Industry benchmarks ───────────────────────────────────────────────────────
INDUSTRY_BENCHMARKS = {
    'saas': """SaaS Industry Benchmarks (reference these when interpreting the data):
- Monthly churn rate: excellent <1%, healthy <2%, concerning >5%
- Net Revenue Retention (NRR): world-class >130%, excellent >110%, good >100%
- CAC Payback Period: excellent <12 months, good <18 months
- Gross Margin: typical 70-85%
- MoM MRR growth: strong >10%, healthy 5-10%
- LTV:CAC ratio: excellent >3:1""",

    'ecommerce': """E-commerce Industry Benchmarks (reference these when interpreting the data):
- Conversion rate: average 1-3%, strong >3%, excellent >5%
- Cart abandonment rate: average 70%, good <65%, excellent <55%
- Customer return rate: good >30%, excellent >45%
- Gross margin: typical 40-60% for physical goods
- Email open rate: average 15-25%
- Average order value growth: target >5% QoQ""",

    'finance': """Finance Industry Benchmarks (reference these when interpreting the data):
- Operating expense ratio: lower is better, varies widely by sub-sector
- Return on Assets (ROA): good >1%, excellent >2%
- Net Interest Margin: typical 2-4% for banks
- Cost-to-Income ratio: efficient <50%, concerning >70%
- Loan-to-Deposit ratio: healthy 80-90%""",

    'marketing': """Marketing Industry Benchmarks (reference these when interpreting the data):
- Email open rate: average 20-25%, strong >30%
- Click-through rate (CTR): average 2-5%, strong >5%
- Lead-to-customer conversion: good >10%, excellent >20%
- Cost per lead (CPL): varies by channel; lower is better
- Marketing ROI: good >5:1, excellent >10:1
- Social media engagement rate: average 1-3%""",

    'operations': """Operations Industry Benchmarks (reference these when interpreting the data):
- On-time delivery rate: good >95%, excellent >98%
- Inventory turnover: higher is generally better (industry-specific)
- Order accuracy rate: good >98%, excellent >99.5%
- First-pass yield: good >95%
- Employee productivity: compare against prior periods for trends""",

    'hr': """HR Industry Benchmarks (reference these when interpreting the data):
- Employee turnover rate: healthy <10%, concerning >20%
- Time-to-hire: good <30 days, excellent <21 days
- Offer acceptance rate: strong >90%
- Employee engagement score: good >70%, excellent >80%
- Absenteeism rate: healthy <2%""",
}

# ── Database ──────────────────────────────────────────────────────────────────
def get_db():
    db = sqlite3.connect(DB_PATH)
    db.row_factory = sqlite3.Row
    return db

def init_db():
    with get_db() as db:
        db.execute('''
            CREATE TABLE IF NOT EXISTS users (
                id                   INTEGER PRIMARY KEY AUTOINCREMENT,
                email                TEXT UNIQUE NOT NULL,
                token                TEXT UNIQUE NOT NULL,
                source               TEXT DEFAULT 'web',
                plan                 TEXT DEFAULT 'free',
                usage_count          INTEGER DEFAULT 0,
                usage_reset_month    TEXT DEFAULT '',
                stripe_customer_id   TEXT DEFAULT '',
                stripe_subscription_id TEXT DEFAULT '',
                created_at           TEXT DEFAULT (datetime('now'))
            )
        ''')
        db.execute('''
            CREATE TABLE IF NOT EXISTS templates (
                id          INTEGER PRIMARY KEY AUTOINCREMENT,
                user_id     INTEGER NOT NULL,
                name        TEXT NOT NULL,
                audience    TEXT DEFAULT 'executive',
                industry    TEXT DEFAULT 'other',
                tone        TEXT DEFAULT 'neutral',
                context     TEXT DEFAULT '',
                outputs     TEXT DEFAULT '["summary","bullets","recommendations","email"]',
                created_at  TEXT DEFAULT (datetime('now')),
                FOREIGN KEY(user_id) REFERENCES users(id)
            )
        ''')
        db.execute('''
            CREATE TABLE IF NOT EXISTS reports (
                id         TEXT PRIMARY KEY,
                user_id    INTEGER,
                title      TEXT DEFAULT 'DataSpeak Report',
                result     TEXT NOT NULL,
                audience   TEXT DEFAULT 'executive',
                industry   TEXT DEFAULT 'other',
                mode       TEXT DEFAULT 'analyse',
                outputs    TEXT DEFAULT '["summary","bullets","recommendations","email"]',
                views      INTEGER DEFAULT 0,
                created_at TEXT DEFAULT (datetime('now'))
            )
        ''')
        cols = [r[1] for r in db.execute("PRAGMA table_info(users)").fetchall()]
        migrations = {
            'plan':                   "ALTER TABLE users ADD COLUMN plan TEXT DEFAULT 'free'",
            'usage_count':            "ALTER TABLE users ADD COLUMN usage_count INTEGER DEFAULT 0",
            'usage_reset_month':      "ALTER TABLE users ADD COLUMN usage_reset_month TEXT DEFAULT ''",
            'stripe_customer_id':     "ALTER TABLE users ADD COLUMN stripe_customer_id TEXT DEFAULT ''",
            'stripe_subscription_id': "ALTER TABLE users ADD COLUMN stripe_subscription_id TEXT DEFAULT ''",
        }
        for col, sql in migrations.items():
            if col not in cols:
                db.execute(sql)
        db.commit()

init_db()

# ── Helpers ───────────────────────────────────────────────────────────────────
def reset_usage_if_needed(db, user):
    this_month = date.today().strftime('%Y-%m')
    if user['usage_reset_month'] != this_month:
        if not user['usage_reset_month']:
            db.execute('UPDATE users SET usage_reset_month=? WHERE id=?',
                       (this_month, user['id']))
            db.commit()
            return user['usage_count']
        db.execute('UPDATE users SET usage_count=0, usage_reset_month=? WHERE id=?',
                   (this_month, user['id']))
        db.commit()
        return 0
    return user['usage_count']

def user_can_generate(user):
    if user['plan'] == 'pro':
        return True, None
    if user['usage_count'] >= FREE_LIMIT:
        return False, f'Free limit reached ({FREE_LIMIT}/month). Upgrade to Pro for unlimited analyses.'
    return True, None

def get_user_by_token(token):
    with get_db() as db:
        return db.execute('SELECT * FROM users WHERE token = ?', (token,)).fetchone()

def render_report_html(report):
    """Render a beautiful public HTML page for a shared report."""
    result = report['result']
    title  = report['title'] or 'DataSpeak Report'

    def badge(text):
        if '[FINDING]' in text:
            t = text.replace('[FINDING]', '<span class="badge finding">FINDING</span>')
        else:
            t = text
        if '[ASSUMPTION]' in text or '[ASSUMPTION]' in t:
            t = t.replace('[ASSUMPTION]', '<span class="badge assumption">ASSUMPTION</span>')
        return t

    section_headers = ['PLAIN ENGLISH SUMMARY','KEY INSIGHTS','RECOMMENDATIONS',
                       'EMAIL FORMAT','COMPARISON SUMMARY','KEY DIFFERENCES',
                       'WHAT IMPROVED','WHAT DECLINED','RED FLAGS','ROOT CAUSES',
                       'RISK LEVEL','IMMEDIATE ACTIONS']

    html_sections = []
    current_h = None
    current_lines = []

    for line in result.split('\n'):
        line = line.rstrip()
        matched = False
        for h in section_headers:
            if line.upper().startswith(h):
                if current_h:
                    html_sections.append((current_h, current_lines))
                current_h = h
                current_lines = []
                matched = True
                break
        if not matched and current_h is not None:
            current_lines.append(line)

    if current_h:
        html_sections.append((current_h, current_lines))

    sections_html = ''
    icons = {
        'PLAIN ENGLISH SUMMARY': '📋',
        'KEY INSIGHTS': '💡',
        'RECOMMENDATIONS': '🎯',
        'EMAIL FORMAT': '✉️',
        'COMPARISON SUMMARY': '🔄',
        'KEY DIFFERENCES': '📊',
        'WHAT IMPROVED': '📈',
        'WHAT DECLINED': '📉',
        'RED FLAGS': '🚨',
        'ROOT CAUSES': '🔍',
        'RISK LEVEL': '⚠️',
        'IMMEDIATE ACTIONS': '⚡',
    }
    for h, lines in html_sections:
        icon = icons.get(h, '📌')
        content_html = ''
        for ln in lines:
            if not ln.strip():
                continue
            if ln.startswith('- '):
                content_html += f'<li>{badge(ln[2:])}</li>'
            else:
                if content_html and not content_html.endswith('</ul>'):
                    pass
                content_html += f'<p>{badge(ln)}</p>'
        # wrap consecutive li in ul
        content_html = content_html.replace('</li><p>', '</li></ul><p>').replace('<p><li>', '<p></p><ul><li>')
        content_html = content_html.replace('<li>', '<ul><li>',1) if '<li>' in content_html and not content_html.startswith('<ul>') else content_html
        if '<li>' in content_html and '</ul>' not in content_html:
            content_html += '</ul>'

        sections_html += f'''
        <div class="section">
          <div class="section-header">{icon} {h.title()}</div>
          <div class="section-body">{content_html}</div>
        </div>'''

    meta_parts = []
    if report['audience'] and report['audience'] != 'executive':
        meta_parts.append(report['audience'].title())
    if report['industry'] and report['industry'] != 'other':
        meta_parts.append(report['industry'].title())
    meta_parts.append(report['created_at'][:10])

    return f'''<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>{title} — DataSpeak</title>
<meta property="og:title" content="{title}">
<meta property="og:description" content="Data analysis report generated by DataSpeak">
<style>
  *{{box-sizing:border-box;margin:0;padding:0}}
  body{{font-family:'Inter',system-ui,sans-serif;background:#0d0f1a;color:#e8eaf0;min-height:100vh;}}
  .topbar{{background:#13162a;border-bottom:1px solid #1e2340;padding:14px 24px;display:flex;align-items:center;gap:12px;}}
  .logo{{font-size:18px;font-weight:800;color:#f7c948;}}
  .logo-sub{{font-size:13px;color:#8aa3c8;margin-left:auto;}}
  .container{{max-width:860px;margin:0 auto;padding:40px 24px 80px;}}
  .report-title{{font-size:32px;font-weight:800;color:#fff;margin-bottom:8px;line-height:1.2;}}
  .report-meta{{font-size:13px;color:#8aa3c8;margin-bottom:32px;}}
  .section{{background:#13162a;border:1px solid #1e2340;border-radius:14px;padding:24px 28px;margin-bottom:20px;}}
  .section-header{{font-size:15px;font-weight:700;color:#f7c948;margin-bottom:14px;letter-spacing:0.03em;}}
  .section-body p{{font-size:15px;line-height:1.7;color:#c8d0e0;margin-bottom:8px;}}
  .section-body ul{{padding-left:20px;}}
  .section-body li{{font-size:15px;line-height:1.7;color:#c8d0e0;margin-bottom:6px;}}
  .badge{{display:inline-block;font-size:10px;font-weight:700;padding:2px 6px;border-radius:4px;margin-right:6px;vertical-align:middle;letter-spacing:0.05em;}}
  .badge.finding{{background:#003d2a;color:#00d48a;border:1px solid #00d48a44;}}
  .badge.assumption{{background:#3d2a00;color:#f7c948;border:1px solid #f7c94844;}}
  .cta-bar{{background:#13162a;border:1px solid #f7c94822;border-radius:14px;padding:24px 28px;text-align:center;margin-top:32px;}}
  .cta-bar p{{color:#8aa3c8;font-size:14px;margin-bottom:14px;}}
  .cta-btn{{display:inline-block;background:linear-gradient(135deg,#f7c948,#e6b800);color:#1a1200;font-weight:800;padding:12px 28px;border-radius:10px;text-decoration:none;font-size:15px;}}
  .views{{font-size:12px;color:#8aa3c8;text-align:right;margin-bottom:16px;}}
</style>
</head>
<body>
<div class="topbar">
  <div class="logo">📊 DataSpeak</div>
  <div class="logo-sub">Shared Report</div>
</div>
<div class="container">
  <div class="views">{report['views']} views</div>
  <div class="report-title">{title}</div>
  <div class="report-meta">{' · '.join(meta_parts)}</div>
  {sections_html}
  <div class="cta-bar">
    <p>This report was generated by DataSpeak — turn any data into clear insights in seconds.</p>
    <a class="cta-btn" href="{APP_URL}">Try DataSpeak Free →</a>
  </div>
</div>
</body>
</html>'''

# ── Routes ────────────────────────────────────────────────────────────────────
@app.route('/')
def index():
    return send_from_directory(app.static_folder, 'index.html')

@app.route('/r/<report_id>')
def view_report(report_id):
    with get_db() as db:
        report = db.execute('SELECT * FROM reports WHERE id=?', (report_id,)).fetchone()
        if not report:
            return '<h1 style="font-family:sans-serif;padding:40px;color:#333">Report not found.</h1>', 404
        db.execute('UPDATE reports SET views=views+1 WHERE id=?', (report_id,))
        db.commit()
    return render_report_html(report)

@app.route('/<path:path>')
def static_files(path):
    return send_from_directory(app.static_folder, path)

@app.route('/api/signup', methods=['POST'])
def signup():
    data   = request.json or {}
    email  = (data.get('email') or '').strip().lower()
    source = data.get('source', 'web')

    if not email or '@' not in email:
        return jsonify({'error': 'Please enter a valid email address.'}), 400

    with get_db() as db:
        existing = db.execute('SELECT * FROM users WHERE email = ?', (email,)).fetchone()
        if existing:
            usage = reset_usage_if_needed(db, existing)
            return jsonify({
                'token': existing['token'], 'email': email, 'returning': True,
                'plan': existing['plan'], 'usage': usage, 'limit': FREE_LIMIT
            })
        token = secrets.token_urlsafe(32)
        db.execute('INSERT INTO users (email, token, source) VALUES (?, ?, ?)',
                   (email, token, source))
        db.commit()

    return jsonify({'token': token, 'email': email, 'returning': False,
                    'plan': 'free', 'usage': 0, 'limit': FREE_LIMIT})

@app.route('/api/status', methods=['POST'])
def status():
    data  = request.json or {}
    token = (data.get('token') or '').strip()
    if not token:
        return jsonify({'error': 'Token required.'}), 400

    with get_db() as db:
        user = db.execute('SELECT * FROM users WHERE token = ?', (token,)).fetchone()
        if not user:
            return jsonify({'error': 'Invalid token.'}), 404
        usage = reset_usage_if_needed(db, user)

    return jsonify({
        'valid': True, 'email': user['email'],
        'plan': user['plan'], 'usage': usage, 'limit': FREE_LIMIT
    })

@app.route('/api/admin/stats')
def admin_stats():
    with get_db() as db:
        total  = db.execute('SELECT COUNT(*) as c FROM users').fetchone()['c']
        today  = db.execute("SELECT COUNT(*) as c FROM users WHERE date(created_at)=date('now')").fetchone()['c']
        pro    = db.execute("SELECT COUNT(*) as c FROM users WHERE plan='pro'").fetchone()['c']
        reports= db.execute('SELECT COUNT(*) as c FROM reports').fetchone()['c']
        recent = db.execute('SELECT email, plan, created_at FROM users ORDER BY created_at DESC LIMIT 10').fetchall()
    return jsonify({
        'total_signups': total, 'today': today, 'pro_users': pro,
        'total_reports': reports,
        'mrr_estimate': pro * 9,
        'recent': [{'email': r['email'], 'plan': r['plan'], 'joined': r['created_at']} for r in recent]
    })

# ── Share report ──────────────────────────────────────────────────────────────
@app.route('/api/share', methods=['POST'])
def share_report():
    data  = request.json or {}
    token = (data.get('token') or '').strip()
    user  = get_user_by_token(token)
    if not user:
        return jsonify({'error': 'Invalid token.'}), 401

    result  = data.get('result', '')
    title   = (data.get('title') or 'DataSpeak Report').strip()
    audience = data.get('audience', 'executive')
    industry = data.get('industry', 'other')
    mode     = data.get('mode', 'analyse')
    outputs  = json.dumps(data.get('outputs', ['summary','bullets','recommendations','email']))

    if not result:
        return jsonify({'error': 'No report to share.'}), 400

    report_id = secrets.token_urlsafe(12)
    with get_db() as db:
        db.execute(
            'INSERT INTO reports (id, user_id, title, result, audience, industry, mode, outputs) VALUES (?,?,?,?,?,?,?,?)',
            (report_id, user['id'], title, result, audience, industry, mode, outputs)
        )
        db.commit()

    return jsonify({'url': f'{APP_URL}/r/{report_id}', 'id': report_id})

# ── Templates ─────────────────────────────────────────────────────────────────
@app.route('/api/templates/save', methods=['POST'])
def save_template():
    data  = request.json or {}
    token = (data.get('token') or '').strip()
    user  = get_user_by_token(token)
    if not user:
        return jsonify({'error': 'Invalid token.'}), 401

    name = (data.get('name') or '').strip()
    if not name:
        return jsonify({'error': 'Template name required.'}), 400

    audience = data.get('audience', 'executive')
    industry = data.get('industry', 'other')
    tone     = data.get('tone', 'neutral')
    context  = data.get('context', '')
    outputs  = json.dumps(data.get('outputs', ['summary','bullets','recommendations','email']))

    with get_db() as db:
        existing = db.execute(
            'SELECT id FROM templates WHERE user_id=? AND name=?', (user['id'], name)
        ).fetchone()
        if existing:
            db.execute(
                'UPDATE templates SET audience=?, industry=?, tone=?, context=?, outputs=? WHERE id=?',
                (audience, industry, tone, context, outputs, existing['id'])
            )
        else:
            db.execute(
                'INSERT INTO templates (user_id, name, audience, industry, tone, context, outputs) VALUES (?,?,?,?,?,?,?)',
                (user['id'], name, audience, industry, tone, context, outputs)
            )
        db.commit()

    return jsonify({'ok': True, 'name': name})

@app.route('/api/templates/list', methods=['POST'])
def list_templates():
    data  = request.json or {}
    token = (data.get('token') or '').strip()
    user  = get_user_by_token(token)
    if not user:
        return jsonify({'error': 'Invalid token.'}), 401

    with get_db() as db:
        rows = db.execute(
            'SELECT * FROM templates WHERE user_id=? ORDER BY created_at DESC LIMIT 20',
            (user['id'],)
        ).fetchall()

    return jsonify({'templates': [{
        'id': r['id'], 'name': r['name'], 'audience': r['audience'],
        'industry': r['industry'], 'tone': r['tone'],
        'context': r['context'], 'outputs': json.loads(r['outputs'])
    } for r in rows]})

@app.route('/api/templates/delete', methods=['POST'])
def delete_template():
    data  = request.json or {}
    token = (data.get('token') or '').strip()
    user  = get_user_by_token(token)
    if not user:
        return jsonify({'error': 'Invalid token.'}), 401

    with get_db() as db:
        db.execute('DELETE FROM templates WHERE id=? AND user_id=?', (data.get('id'), user['id']))
        db.commit()

    return jsonify({'ok': True})

# ── Chat with data ────────────────────────────────────────────────────────────
@app.route('/api/chat', methods=['POST'])
def chat_with_data():
    data  = request.json or {}
    token = (data.get('token') or '').strip()
    user  = get_user_by_token(token)
    if not user:
        return jsonify({'error': 'Invalid token.'}), 401

    original_data = (data.get('data') or '')[:3000]
    analysis      = (data.get('analysis') or '')[:2000]
    question      = (data.get('question') or '').strip()

    if not question:
        return jsonify({'error': 'Question required.'}), 400
    if not ANTHROPIC_KEY:
        return jsonify({'error': 'Server not configured.'}), 500

    prompt = f"""You are a senior data analyst assistant. The user ran an analysis on their data and is asking a follow-up question.

ORIGINAL DATA:
{original_data}

PREVIOUS ANALYSIS:
{analysis}

FOLLOW-UP QUESTION: {question}

Answer directly and concisely. Reference specific numbers from the data. Be analytical, not vague.
Plain text only. No markdown, no asterisks, no bullet dashes unless listing items."""

    try:
        client  = anthropic.Anthropic(api_key=ANTHROPIC_KEY)
        message = client.messages.create(
            model='claude-haiku-4-5-20251001',
            max_tokens=600,
            messages=[{'role': 'user', 'content': prompt}]
        )
        return jsonify({'answer': message.content[0].text})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ── Compare two datasets ──────────────────────────────────────────────────────
@app.route('/api/compare', methods=['POST'])
def compare_datasets():
    data  = request.json or {}
    token = (data.get('token') or '').strip()

    with get_db() as db:
        user = db.execute('SELECT * FROM users WHERE token = ?', (token,)).fetchone()
        if not user:
            return jsonify({'error': 'Invalid token.'}), 401
        reset_usage_if_needed(db, user)
        user = db.execute('SELECT * FROM users WHERE token = ?', (token,)).fetchone()
        can, err = user_can_generate(user)
        if not can:
            return jsonify({'error': err, 'upgrade': True,
                            'usage': user['usage_count'], 'limit': FREE_LIMIT}), 403

    data1    = (data.get('data1') or '').strip()
    data2    = (data.get('data2') or '').strip()
    context  = (data.get('context') or '').strip()
    audience = data.get('audience', 'executive')
    label1   = (data.get('label1') or 'Dataset A').strip()
    label2   = (data.get('label2') or 'Dataset B').strip()

    if not data1 or not data2:
        return jsonify({'error': 'Both datasets are required for comparison.'}), 400
    if not ANTHROPIC_KEY:
        return jsonify({'error': 'Server not configured.'}), 500

    audience_labels = {
        'ceo': 'CEO / Founder', 'cfo': 'CFO / Finance Director',
        'executive': 'Executive / Senior Manager', 'client': 'External Client',
        'team': 'Technical Team', 'board': 'Board of Directors',
        'general': 'General audience', 'professor': 'Professor / Academic Grader',
    }
    audience_label = audience_labels.get(audience, audience)

    prompt = f"""You are an expert data analyst. Compare these two datasets and produce a structured comparison report for a {audience_label}.

{label1}:
{data1}

{label2}:
{data2}

{f'COMPARISON QUESTION: {context}' if context else ''}

Produce the following sections with EXACT headers:

COMPARISON SUMMARY
2-3 sentences summarising the key story between {label1} and {label2}.

KEY DIFFERENCES
4-6 specific differences. Prefix each with [FINDING] if directly in the data or [ASSUMPTION] if inferred.

WHAT IMPROVED
What is better/higher/stronger in {label2} vs {label1}. If nothing, say so.

WHAT DECLINED
What is worse/lower/weaker in {label2} vs {label1}. If nothing, say so.

RECOMMENDATIONS
3-4 specific actions based on this comparison.

Plain text only. No markdown. No asterisks. Section headers in uppercase exactly as shown. Reference actual numbers."""

    try:
        client  = anthropic.Anthropic(api_key=ANTHROPIC_KEY)
        message = client.messages.create(
            model='claude-haiku-4-5-20251001',
            max_tokens=1500,
            messages=[{'role': 'user', 'content': prompt}]
        )
        with get_db() as db:
            db.execute('UPDATE users SET usage_count=usage_count+1 WHERE token=?', (token,))
            db.commit()
            updated = db.execute('SELECT usage_count, plan FROM users WHERE token=?', (token,)).fetchone()

        return jsonify({
            'result': message.content[0].text,
            'usage': updated['usage_count'],
            'limit': FREE_LIMIT,
            'plan': updated['plan']
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ── Export: Word doc ──────────────────────────────────────────────────────────
@app.route('/api/export/word', methods=['POST'])
def export_word():
    data  = request.json or {}
    token = (data.get('token') or '').strip()
    user  = get_user_by_token(token)
    if not user:
        return jsonify({'error': 'Invalid token.'}), 401

    text     = data.get('text', '')
    title    = data.get('title', 'DataSpeak Report')
    audience = data.get('audience', '')
    industry = data.get('industry', '')

    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()

        heading = doc.add_heading(title, 0)
        heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
        if heading.runs:
            heading.runs[0].font.color.rgb = RGBColor(0x1a, 0x56, 0xdb)

        meta_parts = []
        if audience:
            meta_parts.append(f'Audience: {audience}')
        if industry and industry != 'other':
            meta_parts.append(f'Industry: {industry.title()}')
        meta_parts.append(f'Generated: {date.today().strftime("%B %d, %Y")}')
        meta_parts.append('Generated by DataSpeak')

        meta = doc.add_paragraph(' · '.join(meta_parts))
        if meta.runs:
            meta.runs[0].font.size = Pt(9)
            meta.runs[0].font.color.rgb = RGBColor(0x6b, 0x7f, 0xa3)

        doc.add_paragraph()

        section_headers = ['PLAIN ENGLISH SUMMARY','KEY INSIGHTS','RECOMMENDATIONS',
                           'EMAIL FORMAT','COMPARISON SUMMARY','KEY DIFFERENCES',
                           'WHAT IMPROVED','WHAT DECLINED','RED FLAGS','ROOT CAUSES',
                           'RISK LEVEL','IMMEDIATE ACTIONS']
        lines = text.split('\n')

        for line in lines:
            line = line.rstrip()
            if not line:
                continue

            is_header = False
            for h in section_headers:
                if line.upper().startswith(h):
                    p = doc.add_heading(h, level=1)
                    if p.runs:
                        p.runs[0].font.color.rgb = RGBColor(0x1a, 0x56, 0xdb)
                    is_header = True
                    break

            if is_header:
                continue

            if line.startswith('- '):
                content = line[2:].strip()
                if '[FINDING]' in content:
                    content = content.replace('[FINDING]', '').strip()
                    p = doc.add_paragraph(style='List Bullet')
                    run = p.add_run('FINDING  ')
                    run.bold = True
                    run.font.color.rgb = RGBColor(0x00, 0x99, 0x66)
                    p.add_run(content)
                elif '[ASSUMPTION]' in content:
                    content = content.replace('[ASSUMPTION]', '').strip()
                    p = doc.add_paragraph(style='List Bullet')
                    run = p.add_run('ASSUMPTION  ')
                    run.bold = True
                    run.font.color.rgb = RGBColor(0xd2, 0x99, 0x22)
                    p.add_run(content)
                else:
                    p = doc.add_paragraph(style='List Bullet')
                    p.add_run(content)
            else:
                doc.add_paragraph(line)

        doc.add_paragraph()
        footer_p = doc.add_paragraph('Generated by DataSpeak — dataspeak-vydp.onrender.com')
        if footer_p.runs:
            footer_p.runs[0].font.size = Pt(8)
            footer_p.runs[0].font.color.rgb = RGBColor(0x6b, 0x7f, 0xa3)

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)

        filename = f'DataSpeak_Report_{date.today().strftime("%Y%m%d")}.docx'
        return send_file(buf, as_attachment=True, download_name=filename,
                         mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document')

    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ── Export: PowerPoint slides ─────────────────────────────────────────────────
@app.route('/api/export/slides', methods=['POST'])
def export_slides():
    data  = request.json or {}
    token = (data.get('token') or '').strip()
    user  = get_user_by_token(token)
    if not user:
        return jsonify({'error': 'Invalid token.'}), 401

    text     = data.get('text', '')
    title    = data.get('title', 'DataSpeak Report')
    audience = data.get('audience', '')
    industry = data.get('industry', '')

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        BG    = RGBColor(0x0d, 0x0f, 0x1a)
        GOLD  = RGBColor(0xf7, 0xc9, 0x48)
        WHITE = RGBColor(0xff, 0xff, 0xff)
        MUTED = RGBColor(0x8a, 0xa3, 0xc8)
        GREEN = RGBColor(0x00, 0xd4, 0x8a)

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        def bg(slide):
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = BG

        def tb(slide, txt, l, t, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
            box = slide.shapes.add_textbox(l, t, w, h)
            tf  = box.text_frame
            tf.word_wrap = True
            p   = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = txt
            run.font.size  = Pt(size)
            run.font.bold  = bold
            run.font.color.rgb = color

        def gold_bar(slide):
            bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.1), Inches(7.5))
            bar.fill.solid(); bar.fill.fore_color.rgb = GOLD
            bar.line.fill.background()

        def divider(slide, top):
            ln = slide.shapes.add_shape(1, Inches(0.3), top, Inches(12.7), Inches(0.03))
            ln.fill.solid(); ln.fill.fore_color.rgb = GOLD
            ln.line.fill.background()

        # ── Title slide ──────────────────────────────────────────────────────
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        bg(sl)
        gold_bar(sl)
        tb(sl, 'DataSpeak', Inches(0.5), Inches(1.4), Inches(12), Inches(0.6), size=13, color=GOLD)
        tb(sl, title,       Inches(0.5), Inches(2.1), Inches(12), Inches(1.6), size=38, bold=True)
        meta = []
        if audience: meta.append(audience.title())
        if industry and industry != 'other': meta.append(industry.title())
        meta.append(date.today().strftime('%B %d, %Y'))
        tb(sl, '  ·  '.join(meta), Inches(0.5), Inches(3.8), Inches(12), Inches(0.5), size=14, color=MUTED)
        tb(sl, 'dataspeak-vydp.onrender.com', Inches(0.5), Inches(6.9), Inches(12), Inches(0.4), size=10, color=MUTED)

        # ── Parse sections ───────────────────────────────────────────────────
        HEADERS = ['PLAIN ENGLISH SUMMARY','KEY INSIGHTS','RECOMMENDATIONS',
                   'EMAIL FORMAT','COMPARISON SUMMARY','KEY DIFFERENCES',
                   'WHAT IMPROVED','WHAT DECLINED','RED FLAGS','ROOT CAUSES',
                   'RISK LEVEL','IMMEDIATE ACTIONS']
        ICONS   = {
            'PLAIN ENGLISH SUMMARY':'📋','KEY INSIGHTS':'💡','RECOMMENDATIONS':'🎯',
            'EMAIL FORMAT':'✉️','COMPARISON SUMMARY':'🔄','KEY DIFFERENCES':'📊',
            'WHAT IMPROVED':'📈','WHAT DECLINED':'📉','RED FLAGS':'🚨',
            'ROOT CAUSES':'🔍','RISK LEVEL':'⚠️','IMMEDIATE ACTIONS':'⚡',
        }

        sections, cur_h, cur_lines = [], None, []
        for line in text.split('\n'):
            line = line.rstrip()
            matched = False
            for h in HEADERS:
                if line.upper().startswith(h):
                    if cur_h: sections.append((cur_h, cur_lines))
                    cur_h, cur_lines = h, []
                    matched = True; break
            if not matched and cur_h is not None and line:
                cur_lines.append(line)
        if cur_h: sections.append((cur_h, cur_lines))

        for h, lines in sections:
            if not lines: continue
            icon  = ICONS.get(h, '📌')
            label = f'{icon} {h.title()}'
            content = '\n'.join(
                ln.replace('[FINDING]', '✓ ').replace('[ASSUMPTION]', '◆ ')
                for ln in lines if ln.strip()
            )

            sl = prs.slides.add_slide(prs.slide_layouts[6])
            bg(sl); gold_bar(sl)
            tb(sl, label, Inches(0.3), Inches(0.25), Inches(12.7), Inches(0.7), size=26, bold=True)
            divider(sl, Inches(1.05))
            tb(sl, content, Inches(0.3), Inches(1.15), Inches(12.7), Inches(6.1), size=16)
            tb(sl, 'DataSpeak', Inches(11.3), Inches(7.1), Inches(2), Inches(0.35),
               size=9, color=MUTED, align=PP_ALIGN.RIGHT)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        fname = f'DataSpeak_Slides_{date.today().strftime("%Y%m%d")}.pptx'
        return send_file(buf, as_attachment=True, download_name=fname,
                         mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ── Stripe: checkout session ──────────────────────────────────────────────────
@app.route('/api/create-checkout-session', methods=['POST'])
def create_checkout():
    _stripe_key = os.environ.get('STRIPE_SECRET_KEY', '')
    _price_id   = os.environ.get('STRIPE_PRICE_ID', '')
    _app_url    = os.environ.get('APP_URL', '')
    if not _stripe_key:
        return jsonify({'error': 'Payments not configured yet.'}), 500

    import stripe
    stripe.api_key = _stripe_key

    data  = request.json or {}
    token = (data.get('token') or '').strip()

    with get_db() as db:
        user = db.execute('SELECT * FROM users WHERE token = ?', (token,)).fetchone()
    if not user:
        return jsonify({'error': 'Invalid token.'}), 401
    if user['plan'] == 'pro':
        return jsonify({'error': 'Already on Pro plan.'}), 400

    try:
        if user['stripe_customer_id']:
            customer_id = user['stripe_customer_id']
        else:
            customer = stripe.Customer.create(email=user['email'])
            customer_id = customer.id
            with get_db() as db:
                db.execute('UPDATE users SET stripe_customer_id=? WHERE token=?', (customer_id, token))
                db.commit()

        session = stripe.checkout.Session.create(
            customer=customer_id,
            payment_method_types=['card'],
            line_items=[{'price': _price_id, 'quantity': 1}],
            mode='subscription',
            success_url=f'{_app_url}/?upgraded=1&token={token}',
            cancel_url=f'{_app_url}/?cancelled=1',
            subscription_data={'trial_period_days': 7},
        )
        return jsonify({'url': session.url})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ── Stripe: billing portal ────────────────────────────────────────────────────
@app.route('/api/billing-portal', methods=['POST'])
def billing_portal():
    _stripe_key = os.environ.get('STRIPE_SECRET_KEY', '')
    _app_url    = os.environ.get('APP_URL', '')
    if not _stripe_key:
        return jsonify({'error': 'Payments not configured yet.'}), 500

    import stripe
    stripe.api_key = _stripe_key

    data  = request.json or {}
    token = (data.get('token') or '').strip()

    with get_db() as db:
        user = db.execute('SELECT * FROM users WHERE token = ?', (token,)).fetchone()
    if not user or not user['stripe_customer_id']:
        return jsonify({'error': 'No billing account found.'}), 404

    try:
        portal = stripe.billing_portal.Session.create(
            customer=user['stripe_customer_id'],
            return_url=_app_url,
        )
        return jsonify({'url': portal.url})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ── Stripe webhook ────────────────────────────────────────────────────────────
@app.route('/api/stripe-webhook', methods=['POST'])
def stripe_webhook():
    if not STRIPE_SECRET_KEY:
        return 'Not configured', 200

    import stripe
    stripe.api_key = STRIPE_SECRET_KEY

    payload = request.data
    sig     = request.headers.get('Stripe-Signature', '')

    try:
        event = stripe.Webhook.construct_event(payload, sig, STRIPE_WEBHOOK_SECRET)
    except Exception as e:
        return str(e), 400

    etype = event['type']
    obj   = event['data']['object']

    if etype in ('customer.subscription.created', 'customer.subscription.updated'):
        active = obj.get('status') in ('active', 'trialing')
        plan   = 'pro' if active else 'free'
        with get_db() as db:
            db.execute(
                'UPDATE users SET plan=?, stripe_subscription_id=? WHERE stripe_customer_id=?',
                (plan, obj.get('id'), obj.get('customer'))
            )
            db.commit()

    elif etype in ('customer.subscription.deleted', 'invoice.payment_failed'):
        with get_db() as db:
            db.execute('UPDATE users SET plan=? WHERE stripe_customer_id=?',
                       ('free', obj.get('customer')))
            db.commit()

    return jsonify({'received': True})

# ── Generate insights ─────────────────────────────────────────────────────────
@app.route('/api/generate', methods=['POST'])
def generate():
    data  = request.json or {}
    token = (data.get('token') or '').strip()

    if not token:
        return jsonify({'error': 'Please sign in first.'}), 401

    with get_db() as db:
        user = db.execute('SELECT * FROM users WHERE token = ?', (token,)).fetchone()
        if not user:
            return jsonify({'error': 'Invalid token. Please sign in again.'}), 401
        reset_usage_if_needed(db, user)
        user = db.execute('SELECT * FROM users WHERE token = ?', (token,)).fetchone()
        can, err = user_can_generate(user)
        if not can:
            return jsonify({'error': err, 'upgrade': True,
                            'usage': user['usage_count'], 'limit': FREE_LIMIT}), 403

    raw_data    = (data.get('data') or '').strip()
    context     = (data.get('context') or '').strip()
    audience    = data.get('audience', 'executive')
    industry    = data.get('industry', 'other')
    tone        = data.get('tone', 'neutral')
    outputs     = data.get('outputs', ['summary','bullets','recommendations','email'])
    data_source = data.get('data_source', 'sql')
    mode        = data.get('mode', 'analyse')  # 'analyse' | 'red_flags'

    if not raw_data:
        return jsonify({'error': 'No data provided.'}), 400
    if not ANTHROPIC_KEY:
        return jsonify({'error': 'Server not configured. Contact support.'}), 500

    source_context = {
        'sql':    'SQL query results (rows and columns from a database)',
        'r':      'R statistical output (data frames, summary(), lm() results, ANOVA, correlation matrices)',
        'python': 'Python/pandas output (DataFrame, describe(), groupby results)',
        'excel':  'Excel or spreadsheet data (rows, columns, possibly with pivot table output)',
        'other':  'data table or analysis output'
    }.get(data_source, 'data')

    stat_terms = ''
    if data_source == 'r':
        stat_terms = 'Use appropriate statistical language: coefficients, p-values, R-squared, confidence intervals, standard errors where relevant.'
    elif data_source == 'python':
        stat_terms = 'Reference pandas/numpy conventions where applicable.'

    tone_instructions = {
        'urgent':     'Write with urgency. Highlight risks, gaps, and required immediate action. Use direct, assertive language.',
        'reassuring': 'Write with a calm, measured tone. Acknowledge challenges but emphasise positives, progress, and the plan forward.',
        'neutral':    'Write with a balanced, objective, professional tone.',
    }
    tone_instr = tone_instructions.get(tone, tone_instructions['neutral'])

    audience_labels = {
        'ceo':       'CEO / Founder (focus on strategic impact, growth, and big-picture risk)',
        'cfo':       'CFO / Finance Director (focus on financial performance, cost, and ROI)',
        'executive': 'Executive / Senior Manager (professional, high-level, action-oriented)',
        'client':    'External Client (professional, clear, no internal jargon)',
        'team':      'Technical Team (can include data details and technical specifics)',
        'board':     'Board of Directors (formal, concise, strategic risk and opportunity focus)',
        'general':   'General audience (clear, plain English, minimal jargon)',
        'professor': 'Professor / Academic Grader (rigorous, evidence-based, structured)',
    }
    audience_label = audience_labels.get(audience, audience)
    benchmarks_text = INDUSTRY_BENCHMARKS.get(industry, '')

    # ── Red Flags mode ────────────────────────────────────────────────────────
    if mode == 'red_flags':
        prompt = f"""You are a forensic data analyst specialising in risk detection. Your ONLY job is to find problems, anomalies, risks, and warning signs in this data.

DATA ({source_context}):
{raw_data}

{f'CONTEXT: {context}' if context else ''}
AUDIENCE: {audience_label}
{benchmarks_text}
{stat_terms}

Produce the following sections with EXACT headers:

RED FLAGS
List every anomaly, outlier, missing value, suspicious pattern, and warning sign you detect. Be specific — quote exact numbers. Prefix each with [FINDING] if directly in the data, or [ASSUMPTION] if inferred. List at least 5 items.

ROOT CAUSES
For each red flag, suggest the most likely root cause. Prefix with [FINDING] or [ASSUMPTION].

RISK LEVEL
Assign a single overall risk level: CRITICAL / HIGH / MEDIUM / LOW. Justify in 2 sentences with specific evidence from the data.

IMMEDIATE ACTIONS
3-5 specific things the team must do RIGHT NOW to address these issues. Be direct. No vague advice.

Plain text only. No markdown. No asterisks. Headers in uppercase exactly as shown. Reference actual numbers throughout."""

    else:
        # ── Standard analyse mode ─────────────────────────────────────────────
        output_instructions = {
            'summary':         'PLAIN ENGLISH SUMMARY\nWrite 2-3 clear sentences explaining what this data shows. No jargon.',
            'bullets':         'KEY INSIGHTS\nList 5-7 bullet points. For each, prefix with [FINDING] if directly supported by the data, or [ASSUMPTION] if it is a reasonable inference requiring validation.',
            'recommendations': f'RECOMMENDATIONS\nProvide 3-5 specific, actionable recommendations. For each, prefix with [FINDING] if directly supported or [ASSUMPTION] if inferred.',
            'email':           f'EMAIL FORMAT\nWrite a professional email to a {audience_label} sharing these findings. Include: Subject line, greeting, key findings (3-4 sentences), what it means, recommended next steps, and a sign-off.',
        }
        sections = '\n\n'.join(output_instructions[o] for o in outputs if o in output_instructions)

        prompt = f"""You are an expert data analyst and business communication specialist. Analyze the following {source_context} and produce a professional report.

DATA:
{raw_data}

{f'BUSINESS/RESEARCH QUESTION: {context}' if context else ''}
AUDIENCE: {audience_label}
TONE: {tone_instr}
{stat_terms}
{benchmarks_text}

Provide the following outputs using the EXACT section headers shown below:

{sections}

CONFIDENCE SCORING RULES:
- [FINDING] = directly and clearly supported by numbers in the data
- [ASSUMPTION] = reasonable inference not directly stated in the data, requiring further validation

Where industry benchmarks are provided, compare the data against them.

Be specific — reference actual numbers and patterns. Write like a senior analyst presenting to real decision-makers.

FORMATTING: Plain text only. No markdown. No asterisks (*), no bold (**text**), no ## or ### headers. For bullet lists use a dash and space (- item). Write section headers exactly as shown above in plain uppercase."""

    try:
        client  = anthropic.Anthropic(api_key=ANTHROPIC_KEY)
        message = client.messages.create(
            model='claude-haiku-4-5-20251001',
            max_tokens=2000,
            messages=[{'role': 'user', 'content': prompt}]
        )
        with get_db() as db:
            db.execute('UPDATE users SET usage_count=usage_count+1 WHERE token=?', (token,))
            db.commit()
            updated = db.execute('SELECT usage_count, plan FROM users WHERE token=?', (token,)).fetchone()

        return jsonify({
            'result': message.content[0].text,
            'usage':  updated['usage_count'],
            'limit':  FREE_LIMIT,
            'plan':   updated['plan']
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=False)
